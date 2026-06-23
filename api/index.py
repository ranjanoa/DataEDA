import sys
import os

# Fix for PyInstaller --noconsole mode where stdout/stderr are None
# This prevents uvicorn/logging from crashing with AttributeError: 'NoneType' object has no attribute 'isatty'
if sys.stdout is None:
    sys.stdout = open(os.devnull, "w")
if sys.stderr is None:
    sys.stderr = open(os.devnull, "w")

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Body
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
import pandas as pd
import numpy as np
import io
import logging
from typing import List, Optional, Dict

def get_resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
    return os.path.join(base_path, relative_path)

# Add backend directory to path to import data_processor
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from data_processor import (
    load_data, merge_datasets, calculate_stats, downsample_for_plot, 
    calculate_correlations, parse_legend, filter_columns, 
    calculate_derived_var, get_stability_bands, process_extra_dataset,
    get_global_summary, find_shutdown_periods
)

app = FastAPI()

# Configure Templates
templates = Jinja2Templates(directory=get_resource_path(os.path.join("frontend", "templates")))

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount static files
static_dir = get_resource_path(os.path.join("frontend", "static"))
if os.path.exists(static_dir):
    app.mount("/static", StaticFiles(directory=static_dir), name="static")

# Global storage
# In a real multi-user app, this would be a database or session-based storage (e.g. Redis)
current_df: Optional[pd.DataFrame] = None
data_files: Dict[str, pd.DataFrame] = {}
legend_data: Dict[str, Dict] = {}

@app.get("/")
async def read_root():
    # In a real app, use Jinja2 templates. For now, read the file and return.
    pixel_path = get_resource_path(os.path.join("frontend", "templates", "index.html"))
    if os.path.exists(pixel_path):
        with open(pixel_path, 'r', encoding='utf-8') as f:
            return HTMLResponse(content=f.read())
    return {"message": "Backend is running. Frontend index.html not found."}

@app.post("/upload")
async def upload_files(files: List[UploadFile] = File(...)):
    global data_files, legend_data
    uploaded_names = []
    
    for file in files:
        try:
            content = await file.read()
            
            # Check if it's a legend file based on name
            if 'legend' in file.filename.lower() and file.filename.endswith('.csv'):
                legend_data = parse_legend(content)
                uploaded_names.append(f"{file.filename} (Legend)")
            else:
                # Load data file
                df = load_data(content, file.filename)
                data_files[file.filename] = df
                uploaded_names.append(file.filename)
        except Exception as e:
            logging.error(f"Failed to process {file.filename}: {e}")
            return JSONResponse(status_code=400, content={"message": f"Failed to process {file.filename}: {str(e)}"})
            
    return {"message": "Files uploaded successfully", "files": uploaded_names}

@app.post("/upload-legend")
async def upload_legend(file: UploadFile = File(...)):
    global legend_data
    try:
        content = await file.read()
        legend_data = parse_legend(content)
        return {"message": "Legend uploaded successfully", "entries": len(legend_data)}
    except Exception as e:
        logging.error(f"Failed to process legend {file.filename}: {e}")
        return JSONResponse(status_code=400, content={"message": f"Failed to process legend: {str(e)}"})

@app.post("/process")
async def process_data(
    min_nonzero: float = Form(0.0),
    merge_strategy: str = Form('linear')
):
    global current_df, data_files
    
    if not data_files:
        return JSONResponse(status_code=400, content={"message": "No files uploaded"})
    
    try:
        # Sort files? for now take values
        dfs = list(data_files.values())
        
        # Merge
        merged = merge_datasets(dfs, interpolation_method=merge_strategy)
        
        # Filter
        if min_nonzero > 0:
            merged = filter_columns(merged, min_nonzero_pct=min_nonzero)
            
        # Clean column names in merged df
        merged.columns = merged.columns.astype(str).str.replace(' ', '_')
        
        current_df = merged
        
        return {"message": "Data merged successfully", "rows": len(current_df), "columns": list(current_df.columns)}
    except Exception as e:
        logging.error(f"Merge failed: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/add-dataset")
async def add_dataset(
    file: UploadFile = File(...),
    method: str = Form('linear'),
    gap_limit: Optional[float] = Form(None),
    suffix: str = Form('')
):
    global current_df, data_files
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No active dataset to merge into."})
        
    try:
        content = await file.read()
        
        # We don't necessarily add to data_files unless we want to persist it for re-process (which resets current_df).
        # But here we are modifying current_df directly as an additive step.
        # Ideally we should add to data_files and re-merge, but this endpoint implies "add ON TOP of current result".
        # Let's modify current_df.
        
        current_df = process_extra_dataset(
            current_df, 
            content, 
            file.filename, 
            method=method, 
            gap_limit=gap_limit, 
            suffix=suffix
        )
        
        return {"message": f"Dataset {file.filename} added.", "columns": list(current_df.columns)}
    except Exception as e:
        logging.error(f"Add dataset failed: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/add-derived")
async def add_derived_variable(
    name: str = Form(...),
    formula: str = Form(...),
    unit: Optional[str] = Form(None),
    min_val: Optional[str] = Form(None),
    max_val: Optional[str] = Form(None)
):
    global current_df, legend_data
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data active"})
        
    try:
        current_df = calculate_derived_var(current_df, name, formula)
        
        # Update legend metadata if provided
        if unit or min_val or max_val:
            if name not in legend_data:
                legend_data[name] = {}
            
            if unit: legend_data[name]['Unit'] = unit
            if min_val: legend_data[name]['Scale min'] = min_val
            if max_val: legend_data[name]['Scale max'] = max_val
            
        return {"message": f"Variable {name} added", "columns": list(current_df.columns)}
    except Exception as e:
        return JSONResponse(status_code=400, content={"message": str(e)})

@app.get("/columns")
async def get_columns():
    global current_df
    if current_df is None:
        return {"columns": []}
    return {"columns": list(current_df.columns)}

@app.get("/stats")
async def get_stats():
    global current_df
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data processed"})
    
    try:
        stats = calculate_stats(current_df)
        summary = get_global_summary(current_df)
        shutdowns = find_shutdown_periods(current_df)
        return {
            "stats": stats, 
            "summary": summary, 
            "shutdowns": shutdowns
        }
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.get("/legend")
async def get_legend():
    global legend_data
    return {"legend": legend_data}

@app.post("/plot-data")
async def get_plot_data(
    x_var: str = Form(...), 
    y_vars: List[str] = Form(...),
    max_points: int = Form(5000)
):
    global current_df
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data active"})
    
    try:
        # Filter columns
        cols = [x_var] + y_vars
        # Remove duplicates
        cols = list(set(cols))
        
        # Check existence
        missing = [c for c in cols if c not in current_df.columns]
        if missing:
             return JSONResponse(status_code=400, content={"message": f"Columns not found: {missing}"})

        subset = current_df[cols].copy()
        
        # Replace infs/nans with nan first to drop them or clean them
        subset = subset.replace([np.inf, -np.inf], np.nan)
        subset = subset.dropna()
        
        # Downsample
        sampled = downsample_for_plot(subset, max_points)
        
        # Convert to records
        if 'timestamp' in sampled.columns:
            sampled['timestamp'] = sampled['timestamp'].astype(str)
            
        data = sampled.to_dict(orient='list')
        
        # Clean dict of NaNs/Infs for JSON compliance
        cleaned_data = {}
        for col, val_list in data.items():
            cleaned_list = []
            for val in val_list:
                if pd.isna(val) or val == 'nan':
                    cleaned_list.append(None)
                elif isinstance(val, float) and (np.isinf(val) or np.isnan(val)):
                    cleaned_list.append(None)
                else:
                    cleaned_list.append(val)
            cleaned_data[col] = cleaned_list
            
        return {"data": cleaned_data}
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.get("/bands/{var_name}")
async def get_bands(var_name: str):
    global current_df
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data active"})
        
    try:
        bands = get_stability_bands(current_df, var_name)
        return {"bands": bands}
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.get("/heatmap")
async def get_heatmap(method: str = 'pearson'):
    global current_df
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data active"})
    
    try:
        heat = calculate_correlations(current_df, method=method)
        return heat
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/ai-query")
async def ai_query(
    query: str = Form(...),
    context: str = Form(...),
    api_key: Optional[str] = Form(None),
    provider: str = Form("openai"),
    base_url: Optional[str] = Form(None)
):
    import requests
    try:
        if provider == "gemini":
            if not api_key:
                logging.info("No API Key provided. Falling back to local report generation...")
                try:
                    from report_generator import generate_local_report
                    if current_df is not None:
                        content = generate_local_report(current_df, query)
                        return {"choices": [{"message": {"content": content}}]}
                except Exception as local_err:
                    logging.error(f"Local report generation failed: {local_err}")
                return JSONResponse(status_code=400, content={"message": "API Key required for Gemini provider"})
            
            def call_gemini(model):
                target_url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key.strip()}"
                headers = { "Content-Type": "application/json" }
                payload = {
                    "contents": [{
                        "parts": [{
                            "text": f"System: You are a data analysis assistant for kiln process monitoring.\nContext:\n{context}\n\nQuery: {query}"
                        }]
                    }]
                }
                logging.info(f"Sending AI Query to Gemini API using model: {model}")
                return requests.post(target_url, headers=headers, json=payload, timeout=30)

            # Try the default model first (gemini-3-flash is the user's preferred model)
            response = None
            try:
                response = call_gemini("gemini-3-flash")
            except Exception as e:
                logging.error(f"Primary Gemini call failed: {e}")
            
            # If 404/not found, dynamically list models and pick a supported flash/pro model
            if response is None or response.status_code == 404 or "not found" in response.text.lower():
                logging.info("Gemini default model not found. Fetching available models...")
                try:
                    list_url = f"https://generativelanguage.googleapis.com/v1beta/models?key={api_key.strip()}"
                    list_resp = requests.get(list_url, timeout=10)
                    if list_resp.ok:
                        models_data = list_resp.json()
                        models_list = models_data.get("models", [])
                        available_models = []
                        for m in models_list:
                            name = m.get("name", "")
                            methods = m.get("supportedGenerationMethods", [])
                            if name and "generateContent" in methods:
                                available_models.append(name.split("/")[-1])
                        
                        logging.info(f"Available Gemini models supporting generateContent: {available_models}")
                        
                        fallback_model = None
                        # Prefer 3-flash, 3-flash-preview, 3.0-flash, 3.5-flash, 2.5-flash, 2.0-flash, 1.5-flash
                        for pref in ["gemini-3-flash", "gemini-3-flash-preview", "gemini-3.0-flash", "gemini-3.5-flash", "gemini-2.5-flash", "gemini-2.0-flash", "gemini-1.5-flash"]:
                            if pref in available_models:
                                fallback_model = pref
                                break
                        if not fallback_model:
                            for m in available_models:
                                if "flash" in m:
                                    fallback_model = m
                                    break
                        if not fallback_model:
                            for m in available_models:
                                if "pro" in m:
                                    fallback_model = m
                                    break
                        if not fallback_model and available_models:
                            fallback_model = available_models[0]
                            
                        if fallback_model:
                            logging.info(f"Falling back to model: {fallback_model}")
                            response = call_gemini(fallback_model)
                except Exception as list_err:
                    logging.error(f"Error listing Gemini models: {list_err}")
            
            # Secondary fallback with static list of common models if response is still not ok
            if response is None or not response.ok:
                static_fallbacks = ["gemini-3-flash", "gemini-3-flash-preview", "gemini-3.0-flash", "gemini-3.5-flash", "gemini-2.5-flash", "gemini-2.0-flash", "gemini-1.5-flash", "gemini-1.5-flash-latest", "gemini-3.5-pro", "gemini-2.5-pro", "gemini-1.5-pro"]
                for model in static_fallbacks:
                    logging.info(f"Static fallback: Trying {model}")
                    try:
                        fallback_resp = call_gemini(model)
                        if fallback_resp.ok:
                            response = fallback_resp
                            break
                    except Exception as fallback_err:
                        logging.error(f"Fallback model {model} failed with: {fallback_err}")
            
            if response is None or not response.ok:
                status_code = response.status_code if response is not None else 500
                error_text = response.text if response is not None else "Timeout / network error"
                error_msg = f"Gemini Error ({status_code}): {error_text}"
                logging.error(error_msg)
                
                # Perform local fallback on any API error/failure
                logging.info("Gemini call failed. Falling back to local report generation...")
                try:
                    from report_generator import generate_local_report
                    if current_df is not None:
                        content = generate_local_report(current_df, query)
                        return {"choices": [{"message": {"content": content}}]}
                except Exception as local_err:
                    logging.error(f"Local report generation failed: {local_err}")
                return JSONResponse(status_code=status_code, content={"message": error_msg})
            
            res_json = response.json()
            try:
                content = res_json['candidates'][0]['content']['parts'][0]['text']
                # Normalize to match OpenAI format for frontend
                return {"choices": [{"message": {"content": content}}]}
            except (KeyError, IndexError) as e:
                return JSONResponse(status_code=500, content={"message": f"Malformed Gemini response: {str(e)}"})
                
        elif provider == "custom":
            if not base_url:
                return JSONResponse(status_code=400, content={"message": "Base URL required for custom provider"})
            
            headers = { "Content-Type": "application/json" }
            if api_key:
                headers["Authorization"] = f"Bearer {api_key.strip()}"
            
            payload = {
                "model": "gpt-4o", # Assume default or proxy-defined model
                "messages": [
                    {"role": "system", "content": "You are a data analysis assistant."},
                    {"role": "user", "content": f"Context:\n{context}\n\nQuery: {query}"}
                ],
                "max_tokens": 1000
            }
            logging.info(f"Sending AI Query to custom endpoint: {base_url}")
            response = requests.post(base_url, headers=headers, json=payload, timeout=30)
            
            if not response.ok:
                error_msg = f"Custom Provider Error ({response.status_code}): {response.text}"
                logging.error(error_msg)
                return JSONResponse(status_code=response.status_code, content={"message": error_msg})
                
            return response.json()

        else: # Default OpenAI
            if not api_key:
                return JSONResponse(status_code=400, content={"message": "API Key required for OpenAI provider"})
            
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key.strip()}"
            }
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {"role": "system", "content": "You are a data analysis assistant."},
                    {"role": "user", "content": f"Context:\n{context}\n\nQuery: {query}"}
                ],
                "max_tokens": 1000
            }
            logging.info("Sending AI Query to OpenAI API")
            response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload, timeout=30)
            
            if not response.ok:
                error_msg = f"OpenAI Error ({response.status_code}): {response.text}"
                logging.error(error_msg)
                return JSONResponse(status_code=response.status_code, content={"message": error_msg})
                
            return response.json()
            
    except Exception as e:
        logging.error(f"AI Query Exception: {str(e)}", exc_info=True)
        return JSONResponse(status_code=500, content={"message": f"Unexpected backend error: {str(e)}"})

@app.get("/export-data")
async def export_data(format: str = "csv"):
    global current_df
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data active"})
    
    try:
        if format.lower() == "csv":
            stream = io.StringIO()
            current_df.to_csv(stream, index=False)
            response = StreamingResponse(iter([stream.getvalue()]), media_type="text/csv")
            response.headers["Content-Disposition"] = "attachment; filename=export.csv"
            return response
        elif format.lower() == "xlsx":
            stream = io.BytesIO()
            # engine='openpyxl' is default for xlsx usually, but good to specify or let pandas decide
            current_df.to_excel(stream, index=False, engine='openpyxl')
            stream.seek(0)
            response = StreamingResponse(stream, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            response.headers["Content-Disposition"] = "attachment; filename=export.xlsx"
            return response
        else:
            return JSONResponse(status_code=400, content={"message": "Invalid format. Use csv or xlsx."})
    except Exception as e:
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.get("/export-correlations")
async def export_correlations(method: str = 'pearson', format: str = 'csv'):
    global current_df
    if current_df is None:
        return JSONResponse(status_code=400, content={"message": "No data active"})
    
    try:
        # Calculate full correlation matrix for numeric columns
        numeric_df = current_df.select_dtypes(include='number')
        if numeric_df.empty:
            return JSONResponse(status_code=400, content={"message": "No numeric variables found for correlation."})
            
        corr_matrix = numeric_df.corr(method=method)
        
        if format.lower() == "csv":
            stream = io.StringIO()
            corr_matrix.to_csv(stream, index=True) # index=True to include variable names as row labels
            response = StreamingResponse(iter([stream.getvalue()]), media_type="text/csv")
            response.headers["Content-Disposition"] = "attachment; filename=correlation_matrix.csv"
            return response
        elif format.lower() == "xlsx":
            stream = io.BytesIO()
            corr_matrix.to_excel(stream, index=True, engine='openpyxl')
            stream.seek(0)
            response = StreamingResponse(stream, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            response.headers["Content-Disposition"] = "attachment; filename=correlation_matrix.xlsx"
            return response
        else:
            return JSONResponse(status_code=400, content={"message": "Invalid format. Use csv or xlsx."})
    except Exception as e:
        import logging
        logging.error(f"Correlation export failed: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})



@app.post("/generate-report")
async def generate_report(
    data: Dict = Body(...)
):
    """
    Generate an HTML report using the provided analysis data.
    """
    import datetime
    try:
        # Prepare context for template
        context = {
            "title": data.get("title", "Process Analysis"),
            "date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
            "dataset_name": data.get("dataset_name", "Uploaded Data"),
            "summary": data.get("summary", {}),
            "stats": data.get("stats", []),
            "shutdowns": data.get("shutdowns", []),
            "correlations": data.get("correlations", []),
            "ai_summary": data.get("ai_summary", "")
        }
        
        # We need to pass the request object to the template when using Jinja2Templates
        # But here we might just want to return the raw HTML for the browser to open.
        # However, FastAPI's TemplateResponse needs the request.
        # Let's use a dummy request or just render manually if needed.
        from starlette.requests import Request
        # In a real route, the first arg is the request. We can adjust generate_report to take Request.
        
        return JSONResponse(content={"html": templates.get_template("report.html").render(context)})
    except Exception as e:
        logging.error(f"Report generation error: {e}")
        return JSONResponse(status_code=500, content={"message": str(e)})

@app.post("/export-pptx")
async def export_pptx(data: dict = Body(...)):
    import base64
    import io
    import logging
    import re
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    def decode_base64_image(base64_str):
        if ',' in base64_str:
            base64_str = base64_str.split(',')[1]
        image_data = base64.b64decode(base64_str)
        return io.BytesIO(image_data)

    def clean_title(title_str):
        return title_str.strip().replace('**', '').replace('*', '').strip()

    def parse_section_to_blocks(sec_title, raw_content_lines, chart_tag_regex):
        # Split into paragraphs
        paragraphs = []
        current_para = []
        in_table = False
        
        for line in raw_content_lines:
            stripped = line.strip()
            if stripped.startswith('|'):
                if not in_table and current_para:
                    paragraphs.append(("\n".join(current_para), "text"))
                    current_para = []
                in_table = True
                current_para.append(line)
            else:
                if in_table:
                    if current_para:
                        paragraphs.append(("\n".join(current_para), "table"))
                        current_para = []
                    in_table = False
                
                if not stripped:
                    if current_para:
                        paragraphs.append(("\n".join(current_para), "text"))
                        current_para = []
                else:
                    current_para.append(line)
                    
        if current_para:
            paragraphs.append(("\n".join(current_para), "table" if in_table else "text"))

        # Classify
        classified = []
        for text, p_type in paragraphs:
            if p_type == "table":
                classified.append({"text": text, "type": "table", "consumed": False})
            elif re.search(chart_tag_regex, text):
                classified.append({"text": text, "type": "chart_tag", "consumed": False})
            else:
                classified.append({"text": text, "type": "text", "consumed": False})
                
        # Group Chart blocks
        blocks = []
        n = len(classified)
        for idx in range(n):
            item = classified[idx]
            if item["type"] == "chart_tag":
                title = ""
                tag_line = item["text"]
                explanation = ""
                
                # Split item["text"] if it contains multiple lines (e.g., title + tag)
                lines_in_item = [l.strip() for l in item["text"].split('\n') if l.strip()]
                if len(lines_in_item) > 1:
                    title_lines = []
                    for l in lines_in_item:
                        if re.search(chart_tag_regex, l):
                            tag_line = l
                        else:
                            title_lines.append(l)
                    if title_lines:
                        title = "\n".join(title_lines)
                
                # Find title (look back for preceding unconsumed text)
                if not title and idx > 0:
                    prev_item = classified[idx - 1]
                    if prev_item["type"] == "text" and not prev_item["consumed"]:
                        prev_text = prev_item["text"]
                        if len(prev_text) <= 150 or prev_text.strip().startswith('**') or "Graph" in prev_text or "Plot" in prev_text or "Chart" in prev_text:
                            title = prev_text
                            prev_item["consumed"] = True
                            
                # Find explanation (look forward for next unconsumed text)
                if idx + 1 < n:
                    next_item = classified[idx + 1]
                    if next_item["type"] == "text" and not next_item["consumed"]:
                        explanation = next_item["text"]
                        next_item["consumed"] = True
                        
                item["consumed"] = True
                blocks.append({
                    "type": "chart",
                    "tag": tag_line,
                    "title": title,
                    "explanation": explanation,
                    "original_idx": idx
                })
                
        # Now build the final list of blocks, inserting table/text blocks in their original order
        final_blocks = []
        current_text_block = []
        
        for idx in range(n):
            item = classified[idx]
            if item["consumed"]:
                # If we accumulated text, push it first
                if current_text_block:
                    final_blocks.append({
                        "type": "text",
                        "content": "\n\n".join(current_text_block)
                    })
                    current_text_block = []
                    
                # If this was a chart_tag, find the corresponding chart block we created
                if item["type"] == "chart_tag":
                    for cb in blocks:
                        if cb["original_idx"] == idx:
                            final_blocks.append(cb)
                            break
            else:
                if item["type"] == "table":
                    if current_text_block:
                        final_blocks.append({
                            "type": "text",
                            "content": "\n\n".join(current_text_block)
                        })
                        current_text_block = []
                    final_blocks.append({
                        "type": "table",
                        "content": item["text"]
                    })
                else: # unconsumed text
                    current_text_block.append(item["text"])
                    
        if current_text_block:
            final_blocks.append({
                "type": "text",
                "content": "\n\n".join(current_text_block)
            })
            
        return final_blocks

    def extract_tag_features(tag_str):
        tag_match = re.search(r'\[([A-Z0-9]+):[^\]]+\]', tag_str, re.IGNORECASE)
        if tag_match:
            chart_type = tag_match.group(1).lower()
            content = tag_match.group(0).lower()
        else:
            chart_type = ""
            content = tag_str.lower()
            
        tokens = set(re.findall(r'[a-zA-Z0-9_#\-\.]+', content))
        keywords_to_remove = {'x', 'y', 'z', 'color', 'scale', 'jet', 'rdbu', 'hot', 'viridis', 'accent', 'gray'}
        tokens = tokens - keywords_to_remove
        return {
            "type": chart_type,
            "tokens": tokens
        }

    def extract_title_features(title_str):
        title_str = title_str.strip().lower()
        chart_type = ""
        if "scatter 3d" in title_str or "3d scatter" in title_str:
            chart_type = "scatter3d"
        elif "scatter" in title_str:
            chart_type = "scatter"
        elif "dual" in title_str or "axis" in title_str:
            chart_type = "dualplot"
        elif "parallel" in title_str:
            chart_type = "parallel"
        elif "box" in title_str:
            chart_type = "box"
        elif "distribution" in title_str or "histogram" in title_str:
            chart_type = "histogram"
        elif "plot" in title_str:
            chart_type = "plot"
        
        tokens = set(re.findall(r'[a-zA-Z0-9_#\-\.]+', title_str))
        stop_words = {'scatter', '3d', 'dual', 'y-axis', 'axis', 'parallel', 'coordinates', 'box', 'plot', 'distribution', 'histogram', 'color', 'vs'}
        tokens = tokens - stop_words
        return {
            "type": chart_type,
            "tokens": tokens
        }

    def find_best_chart(tag_str, ai_charts, used_indices):
        tag_feats = extract_tag_features(tag_str)
        tag_type = tag_feats["type"]
        tag_tokens = tag_feats["tokens"]
        
        best_idx = -1
        best_score = -1
        
        for idx, chart in enumerate(ai_charts):
            if idx in used_indices:
                continue
                
            title = chart.get("title", "")
            title_feats = extract_title_features(title)
            title_type = title_feats["type"]
            title_tokens = title_feats["tokens"]
            
            type_match = (tag_type == title_type) or (tag_type == "dualplot" and title_type == "timeseries") or (tag_type == "plot" and title_type == "timeseries")
            overlap = len(tag_tokens.intersection(title_tokens))
            score = overlap + (5 if type_match else 0)
            
            if score > best_score:
                best_score = score
                best_idx = idx
                
        if best_idx != -1 and best_score >= 1:
            used_indices.add(best_idx)
            return ai_charts[best_idx]
            
        for idx, chart in enumerate(ai_charts):
            if idx not in used_indices:
                used_indices.add(idx)
                return chart
        return None

    try:
        prs = Presentation()
        # Set slide dimensions to widescreen (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Slide Layout 6 is blank
        blank_layout = prs.slide_layouts[6]
        
        # 1. Title Slide
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 23, 42)
        
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = data.get("title", "Process Stability Analysis Report")
        p.font.name = 'Arial'
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(6, 182, 212) # Cyan Accent
        
        p2 = tf.add_paragraph()
        p2.text = f"\nDataset: {data.get('dataset_name', 'Uploaded Data')}\nDate: {data.get('date', '')}"
        p2.font.name = 'Arial'
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(148, 163, 184) # Slate gray

        # 2. Parse Markdown Sections
        summary_text = data.get("ai_summary", "")
        ai_charts = data.get("ai_charts", [])
        
        lines = summary_text.split('\n')
        sections = []
        current_section = {
            "title": "",
            "content": []
        }
        
        for line in lines:
            stripped = line.strip()
            # Detect headers
            if stripped.startswith('###') or (stripped.startswith('##') and 'STEP' in stripped):
                if current_section["content"] or current_section["title"]:
                    sections.append(current_section)
                title = stripped.lstrip('#').strip().replace('**', '').replace('*', '')
                current_section = {
                    "title": title,
                    "content": []
                }
            else:
                current_section["content"].append(line)
                
        if current_section["content"] or current_section["title"]:
            sections.append(current_section)

        # RegEx to match any Plotly chart tag (case-insensitive)
        chart_tag_regex = r'(?i)\[(?:SCATTER|SCATTER3D|DUALPLOT|PARALLEL|BOX|HISTOGRAM|PLOT):[^\]]+\]'
        used_chart_indices = set()

        # Helper to populate bullet points to a text frame
        def add_bullets_to_tf(tf, bullets):
            tf.word_wrap = True
            first = True
            for b in bullets:
                clean_b = b.strip().replace('**', '').replace('*', '').strip()
                if not clean_b:
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                
                # Format bullets nicely
                if b.strip().startswith('-') or b.strip().startswith('*'):
                    p.text = "• " + clean_b.lstrip('-*').strip()
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(241, 245, 249)
                    p.space_before = Pt(4)
                else:
                    p.text = clean_b
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(241, 245, 249)
                    p.space_before = Pt(6)
                p.font.name = 'Arial'

        for sec in sections:
            blocks = parse_section_to_blocks(sec["title"], sec["content"], chart_tag_regex)
            
            for block in blocks:
                if block["type"] == "table":
                    # Parse Table rows
                    table_lines = [l for l in block["content"].split('\n') if l.strip().startswith('|')]
                    table_data = []
                    for line in table_lines:
                        parts = [p.strip() for p in line.split('|')]
                        if parts and parts[0] == '': parts = parts[1:]
                        if parts and parts[-1] == '': parts = parts[:-1]
                        if parts and all(c == '-' or c == ':' or c == ' ' for p in parts for c in p):
                            continue
                        table_data.append(parts)
                    
                    rows_count = len(table_data)
                    cols_count = len(table_data[0]) if rows_count > 0 else 0
                    
                    if rows_count > 0 and cols_count > 0:
                        slide = prs.slides.add_slide(blank_layout)
                        fill = slide.background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(15, 23, 42)
                        
                        # Section Title
                        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(0.8))
                        tf = title_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = sec["title"] or "Data Summary"
                        p.font.name = 'Arial'
                        p.font.size = Pt(24)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(168, 85, 247) # Purple Accent
                        
                        # Add Native Slide Table
                        left = Inches(1.0)
                        top = Inches(1.4)
                        width = Inches(11.333)
                        height = Inches(min(5.2, 0.35 * rows_count))
                        
                        table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
                        table = table_shape.table
                        
                        for r_idx, row_vals in enumerate(table_data):
                            for c_idx, val in enumerate(row_vals):
                                if r_idx < len(table.rows) and c_idx < len(table.columns):
                                    cell = table.cell(r_idx, c_idx)
                                    cell.text = val
                                    
                                    # Theme cell fills
                                    cell.fill.solid()
                                    if r_idx == 0:
                                        cell.fill.fore_color.rgb = RGBColor(30, 41, 59) # Dark Slate Header
                                    else:
                                        cell.fill.fore_color.rgb = RGBColor(15, 23, 42) # Body dark
                                        
                                    for paragraph in cell.text_frame.paragraphs:
                                        paragraph.font.name = 'Arial'
                                        paragraph.font.size = Pt(10)
                                        paragraph.font.color.rgb = RGBColor(241, 245, 249)
                                        if r_idx == 0:
                                            paragraph.font.bold = True
                                            paragraph.font.color.rgb = RGBColor(6, 182, 212) # Cyan header font
                                            
                elif block["type"] == "chart":
                    chart = find_best_chart(block["tag"], ai_charts, used_chart_indices)
                    img_b64 = chart.get("image_base64", "") if chart else ""
                    
                    if img_b64:
                        slide = prs.slides.add_slide(blank_layout)
                        fill = slide.background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(15, 23, 42)
                        
                        # Section Title
                        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
                        tf = title_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = clean_title(block["title"]) if block["title"] else (sec["title"] or "Process Visualization")
                        p.font.name = 'Arial'
                        p.font.size = Pt(24)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(6, 182, 212) # Cyan accent
                        
                        # Left Column: Image
                        try:
                            img_stream = decode_base64_image(img_b64)
                            slide.shapes.add_picture(img_stream, Inches(0.6), Inches(1.3), width=Inches(6.2), height=Inches(5.4))
                        except Exception as img_err:
                            logging.error(f"Failed to add chart image to slide: {img_err}")
                            err_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(6.2), Inches(2.0))
                            err_tf = err_box.text_frame
                            err_p = err_tf.paragraphs[0]
                            err_p.text = f"[Chart Image Render Error: {str(img_err)}]"
                            err_p.font.color.rgb = RGBColor(239, 68, 68)
                            err_p.font.size = Pt(14)
                            
                        # Right Column: Explanation
                        text_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.733), Inches(5.4))
                        tf_exp = text_box.text_frame
                        
                        exp_bullets = []
                        if block["explanation"]:
                            exp_bullets = [b.strip() for b in block["explanation"].split('\n') if b.strip()]
                        
                        add_bullets_to_tf(tf_exp, exp_bullets)
                    else:
                        # Fallback: Render as text slide if image not found
                        slide = prs.slides.add_slide(blank_layout)
                        fill = slide.background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(15, 23, 42)
                        
                        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8))
                        tf = title_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = clean_title(block["title"]) if block["title"] else (sec["title"] or "Process Insights")
                        p.font.name = 'Arial'
                        p.font.size = Pt(26)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(168, 85, 247) # Purple accent
                        
                        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.2))
                        tf_content = content_box.text_frame
                        
                        bullets = []
                        if block["explanation"]:
                            bullets = [b.strip() for b in block["explanation"].split('\n') if b.strip()]
                        add_bullets_to_tf(tf_content, bullets)
                        
                elif block["type"] == "text":
                    slide = prs.slides.add_slide(blank_layout)
                    fill = slide.background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(15, 23, 42)
                    
                    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8))
                    tf = title_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = sec["title"] or "Executive Insights"
                    p.font.name = 'Arial'
                    p.font.size = Pt(26)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(168, 85, 247) # Purple accent
                    
                    # Split lines into bullet points
                    bullets = [b.strip() for b in block["content"].split('\n') if b.strip()]
                    
                    if len(bullets) > 6:
                        mid = (len(bullets) + 1) // 2
                        bullets_left = bullets[:mid]
                        bullets_right = bullets[mid:]
                        
                        left_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.4), Inches(5.2))
                        add_bullets_to_tf(left_box.text_frame, bullets_left)
                        
                        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.533), Inches(5.2))
                        add_bullets_to_tf(right_box.text_frame, bullets_right)
                    else:
                        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.2))
                        add_bullets_to_tf(content_box.text_frame, bullets)

        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        response = StreamingResponse(ppt_stream, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        response.headers["Content-Disposition"] = "attachment; filename=DataEDA_Report.pptx"
        return response
    except Exception as e:
        logging.error(f"PPTX export failed: {e}", exc_info=True)
        return JSONResponse(status_code=500, content={"message": f"PowerPoint export failed: {str(e)}"})

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)
